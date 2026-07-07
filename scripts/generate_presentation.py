# -*- coding: utf-8 -*-
"""Generate defense presentation for StayJourney thesis."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "presentation_assets"
OUTPUT = ROOT / "Bao_ve_DoAn_StayJourney.pptx"

# Brand colors
PRIMARY = RGBColor(0x1A, 0x3A, 0x5C)   # navy
ACCENT = RGBColor(0xA0, 0x82, 0x6D)    # warm brown (StayJourney)
TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MUTED = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_bar(slide, prs):
    bar = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, Inches(0.12)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, prs, text="StayJourney — Đồ án tốt nghiệp 2026"):
    box = slide.shapes.add_textbox(
        Inches(0.5), prs.slide_height - Inches(0.45),
        prs.slide_width - Inches(1), Inches(0.3)
    )
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_title_block(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(12.2), Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(4)


def add_bullets(slide, items, left=0.55, top=1.35, width=5.8, height=5.5, font_size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
    return box


def add_image(slide, filename, left, top, width, height=None):
    path = ASSETS / filename
    if not path.exists():
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(1))
        box.text_frame.text = f"[Thiếu ảnh: {filename}]"
        return
    if height:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    else:
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    # accent strip
    strip = slide.shapes.add_shape(1, 0, Inches(3.2), prs.slide_width, Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "THIẾT KẾ XÂY DỰNG HỆ THỐNG\nĐẶT PHÒNG VÀ QUẢN LÝ KHÁCH SẠN TRỰC TUYẾN"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11), Inches(2.5))
    tf2 = sub.text_frame
    lines = [
        ("Hệ thống: ", False),
        ("StayJourney", True),
        ("", False),
        ("Sinh viên thực hiện: ", False),
        ("Đoàn Nhật Quang — MSSV 20225911", True),
        ("quang.dn225911@sis.hust.edu.vn", False),
        ("", False),
        ("Giảng viên hướng dẫn: ", False),
        ("ThS. Nguyễn Hồng Phương", True),
        ("", False),
        ("Chương trình: Công nghệ thông tin Việt-Nhật", False),
        ("Khoa Công nghệ Thông tin và Truyền thông — ĐHBK Hà Nội", False),
        ("Hà Nội, 06/2026", False),
    ]
    first = True
    for text, bold in lines:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        p.text = text
        p.font.size = Pt(16 if bold else 14)
        p.font.bold = bold
        p.font.color.rgb = WHITE if bold else RGBColor(0xCC, 0xD6, 0xE0)
        p.space_after = Pt(2)


def slide_outline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "Nội dung trình bày")
    items = [
        "1. Tổng quan đề tài",
        "2. Khảo sát & Phân tích yêu cầu",
        "3. Thiết kế hệ thống",
        "4. Kết quả cài đặt",
        "5. Đóng góp & Giải pháp nổi bật",
        "6. Kết luận & Hướng phát triển",
    ]
    add_bullets(slide, items, font_size=22, top=1.6)
    add_footer(slide, prs)


def slide_overview_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "1. Tổng quan — Giới thiệu đề tài", "Mục 1.1 – 1.2")
    items = [
        "Lý do: Nhu cầu đặt phòng trực tuyến tăng; KS vừa & nhỏ vẫn quản lý thủ công",
        "Hệ quả: Dữ liệu lệch, đặt trùng phòng, đối soát thanh toán khó",
        "Khả thi: Web full-stack (React + Node.js + MongoDB), chi phí phù hợp",
        "Giải pháp: StayJourney — nền tảng thống nhất cho chuỗi KS Việt Nam",
        "4 vai trò: Guest · Owner · Staff · Admin",
        "Mục tiêu: Gom luồng tìm kiếm → đặt phòng → thanh toán → vận hành",
    ]
    add_bullets(slide, items, width=6.2)
    add_image(slide, "usecase_overview.png", 6.9, 1.3, 5.8)
    add_footer(slide, prs)


def slide_survey(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "1. Tổng quan — Khảo sát hiện trạng", "Mục 2.1")
    headers = slide.shapes.add_table(4, 3, Inches(0.55), Inches(1.35), Inches(12.0), Inches(3.2)).table
    data = [
        ("Giải pháp", "Ưu điểm", "Hạn chế / Hướng kế thừa"),
        ("OTA", "Tiếp cận khách lớn, thanh toán tập trung", "Phụ thuộc sàn; khó tùy biến QR, hoàn tiền"),
        ("PMS thương mại", "Vận hành nội bộ sâu", "Chi phí cao; đào tạo phức tạp"),
        ("Website riêng", "Chủ động thương hiệu", "Thiếu quản trị đa vai trò, đa KS"),
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = headers.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11 if r else 12)
                p.font.bold = r == 0
                p.font.color.rgb = PRIMARY if r == 0 else TEXT
    items = [
        "StayJourney kế thừa đặt phòng trực tuyến + vận hành nội bộ gọn nhẹ",
        "Phân quyền theo vai trò & khách sạn; thanh toán QR/VNPay theo từng KS",
    ]
    add_bullets(slide, items, top=4.75, font_size=14, height=1.5)
    add_footer(slide, prs)


def slide_usecase_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "2. Phân tích — Use case tổng quát", "Mục 2.2.1")
    add_image(slide, "usecase_overview.png", 0.7, 1.25, 11.8)
    add_footer(slide, prs)


def slide_usecase_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "2. Phân tích — Phân rã Use case", "Đặt phòng · Quản lý đơn")
    add_image(slide, "usecase_booking.png", 0.4, 1.2, 6.2)
    add_image(slide, "usecase_manage_booking.png", 6.8, 1.2, 6.0)
    add_footer(slide, prs)


def slide_activity_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "2. Phân tích — Quy trình nghiệp vụ", "Đặt phòng & Lưu trú")
    add_image(slide, "flow_booking.png", 0.4, 1.2, 6.2)
    add_image(slide, "flow_stay.png", 6.8, 1.2, 6.0)
    add_footer(slide, prs)


def slide_sequence(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "2. Phân tích — Biểu đồ trình tự", "Đặt phòng online · Hủy đơn")
    add_image(slide, "sequence_online_booking.png", 0.4, 1.15, 6.2)
    add_image(slide, "sequence_cancel_booking.png", 6.8, 1.15, 6.0)
    add_footer(slide, prs)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "3. Thiết kế — Kiến trúc tổng thể", "Three-layer · Client–Server")
    items = [
        "Presentation: React 18 + Vite + SCSS (SPA responsive)",
        "Business: Node.js + Express — booking, payment, pricing…",
        "Data: MongoDB (Mongoose) + Cloudinary",
        "Bảo mật: JWT HttpOnly cookie · CSRF · 2FA OTP",
        "Realtime: Socket.IO · Cron: hủy đơn pending, sale hết hạn",
    ]
    add_bullets(slide, items, width=4.8, font_size=14, top=1.4)
    add_image(slide, "arch_three_layer.png", 5.6, 1.15, 7.0)
    add_footer(slide, prs)


def slide_package_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "3. Thiết kế — Sơ đồ gói & Đặt phòng", "Mục 4.1.3")
    add_image(slide, "package_overview.png", 0.35, 1.15, 6.3)
    add_image(slide, "package_booking_payment.png", 6.75, 1.15, 6.1)
    add_footer(slide, prs)


def slide_class_diagram(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "3. Thiết kế — Biểu đồ lớp", "4 lớp: Boundary · Control · Entity · Infrastructure")
    add_image(slide, "class_diagram.png", 0.6, 1.15, 11.8)
    add_footer(slide, prs)


def slide_erd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "3. Thiết kế — Cơ sở dữ liệu", "ERD logic · 9 collections MongoDB")
    add_image(slide, "erd.png", 0.6, 1.15, 8.5)
    items = [
        "users · hotels · rooms · bookings",
        "reviews · salePromotions · notifications",
        "paymentTransactions · contactMessages",
        "ObjectId tham chiếu · schema Mongoose",
    ]
    add_bullets(slide, items, left=9.4, top=1.8, width=3.4, font_size=13)
    add_footer(slide, prs)


def slide_ui_design(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "3. Thiết kế — Giao diện người dùng", "Responsive · Material Design · SCSS")
    items = [
        "4 portal theo vai trò: / · /owner · /staff · /admin",
        "Bảng màu chủ đạo: #A0826D (nâu ấm) + navy #1A3A5C",
        "Font Gilroy · Toastify phản hồi góc phải trên",
        "Breakpoint: mobile < 768px · tablet · desktop",
        "Component tái sử dụng: modal, calendar, room map",
        "Protected route + redirect theo role sau đăng nhập",
    ]
    add_bullets(slide, items, width=6.0, font_size=14, top=1.45)
    add_image(slide, "ui_guest.png", 6.7, 1.15, 5.9)
    add_footer(slide, prs)


def slide_ui_guest(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "4. Kết quả cài đặt — Giao diện Guest", "Tìm KS · Đặt phòng · Thanh toán")
    add_image(slide, "ui_guest.png", 0.6, 1.15, 11.8)
    add_footer(slide, prs)


def slide_ui_owner(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "4. Kết quả cài đặt — Owner & Thanh toán", "Xác nhận QR · Quản lý đơn")
    add_image(slide, "ui_payment_owner.png", 0.6, 1.15, 11.8)
    add_footer(slide, prs)


def slide_ui_admin(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "4. Kết quả cài đặt — Admin & Kiểm thử", "~65.000 LOC · 161 test cases · 100% pass")
    add_image(slide, "ui_admin.png", 0.5, 1.1, 7.8)
    items = [
        "Jest + Supertest + mongodb-memory-server",
        "Auth, booking lifecycle, payment, authorization",
        "Dev: localhost:3000 / :8001",
        "VNPay Sandbox · Cloudinary · Gmail SMTP",
    ]
    add_bullets(slide, items, left=8.5, top=1.6, width=4.2, font_size=14)
    add_footer(slide, prs)


def slide_contribution_pricing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "5. Đóng góp — Gợi ý giá theo quy tắc", "Rule-based price suggestion")
    items = [
        "3 hệ số: lấp đầy (occ) · mùa/lễ (season) · thứ trong tuần (history 84 ngày)",
        "Công thức: suggested = round(avgPrice × m_occ × m_sea × m_hist)",
        "Kẹp an toàn: 75% – 135% giá trung bình hiện tại",
        "Có giải thích từng hệ số — chỉ áp dụng khi Owner xác nhận",
        "Không dùng ML — phù hợp KS vừa & nhỏ, dữ liệu còn ít",
    ]
    add_bullets(slide, items, width=12, font_size=15, top=1.5)
    add_footer(slide, prs)


def slide_contribution_tech(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "5. Đóng góp — Giải pháp kỹ thuật nổi bật", "Mục 5 & Chương 4")
    items = [
        "Two-channel notification: cá nhân + theo khách sạn (write-before-push)",
        "Socket.IO realtime — Staff/Owner phản hồi nhanh đơn mới",
        "Chặn đặt trùng phòng khi nhiều khách đặt song song",
        "Thời hạn giữ phòng pending + cron tự hủy (Asia/Ho_Chi_Minh)",
        "Đa vai trò trên một nền tảng: refund policy & QR theo từng KS",
        "CSRF double-submit + JWT HttpOnly + 2FA OTP (Admin/Owner)",
    ]
    add_bullets(slide, items, width=12, font_size=15, top=1.5)
    add_footer(slide, prs)


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_top_bar(slide, prs)
    add_title_block(slide, "6. Kết luận", "Đạt được · Hạn chế · Hướng phát triển")
    col1 = [
        "✓ Hoàn thiện StayJourney — 4 vai trò, luồng nghiệp vụ chính",
        "✓ 161/161 test API pass · Giao diện responsive",
        "✓ Tích hợp VNPay Sandbox, QR, Socket.IO, gợi ý giá",
    ]
    col2 = [
        "△ Chưa deploy production / HTTPS / CI-CD",
        "△ VNPay Sandbox; chưa E2E test tự động UI",
        "△ Check-in đơn giản; chưa Web Push, TOTP",
    ]
    col3 = [
        "→ Triển khai cloud + domain thật",
        "→ E2E test (Cypress/Playwright)",
        "→ Mở rộng: đa phòng/đơn, RM nâng cao, mobile app",
    ]
    for i, (title, items) in enumerate([
        ("Đạt được", col1),
        ("Hạn chế", col2),
        ("Hướng PT", col3),
    ]):
        left = 0.55 + i * 4.15
        box = slide.shapes.add_textbox(Inches(left), Inches(1.4), Inches(3.9), Inches(4.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = TEXT
            p.space_after = Pt(6)
    add_footer(slide, prs)


def slide_thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    strip = slide.shapes.add_shape(1, 0, Inches(3.45), prs.slide_width, Inches(0.08))
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CẢM ƠN QUÝ THẦY CÔ\nĐÃ LẮNG NGHE!"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.2))
    tf2 = sub.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Sẵn sàng trả lời câu hỏi"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xCC, 0xD6, 0xE0)
    p2.alignment = PP_ALIGN.CENTER

    info = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8))
    tf3 = info.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "Đoàn Nhật Quang · quang.dn225911@sis.hust.edu.vn · StayJourney"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
    p3.alignment = PP_ALIGN.CENTER


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_cover,
        slide_outline,
        slide_overview_intro,
        slide_survey,
        slide_usecase_overview,
        slide_usecase_detail,
        slide_activity_flow,
        slide_sequence,
        slide_architecture,
        slide_package_design,
        slide_class_diagram,
        slide_erd,
        slide_ui_design,
        slide_ui_guest,
        slide_ui_owner,
        slide_ui_admin,
        slide_contribution_pricing,
        slide_contribution_tech,
        slide_conclusion,
        slide_thanks,
    ]

    for fn in builders:
        fn(prs)

    prs.save(str(OUTPUT))
    print(f"Created: {OUTPUT} ({len(builders)} slides)")


if __name__ == "__main__":
    main()
